#!/usr/bin/env python3
"""vs_pptx.py — PPTX 结构无损导出（python-pptx）。

形状位置（EMU→pt）、填充色、字体、文本、图片引用 —— 全部来自文件本身，零渲染、零丢失。

用法:
  vs_pptx.py --file PATH [--max-shapes 200] [--slide N]
"""
import argparse
import json
import vs_schema as S
import sys
from io import BytesIO


def color_of(fill) -> str | None:
    try:
        if fill is None or fill.type is None:
            return None
        from pptx.enum.dml import MSO_FILL  # type: ignore[import-not-found]

        if fill.type == MSO_FILL.SOLID:
            rgb = fill.fore_color.rgb
            return f"#{rgb}"
        return str(fill.type)
    except Exception:
        return None


def shape_dump(shape, depth: int, out: list):
    emu_pt = 12700.0
    item = {
        "id": shape.shape_id,
        "name": shape.name[:80],
        "type": str(shape.shape_type),
        "pos_pt": [round(shape.left / emu_pt, 2) if shape.left is not None else None,
                   round(shape.top / emu_pt, 2) if shape.top is not None else None],
        "size_pt": [round(shape.width / emu_pt, 2) if shape.width is not None else None,
                    round(shape.height / emu_pt, 2) if shape.height is not None else None],
        "rotation": round(shape.rotation, 1) if shape.rotation else 0,
    }
    if shape.has_text_frame:
        texts = []
        for para in shape.text_frame.paragraphs:
            runs = []
            for run in para.runs:
                fs = run.font.size.pt if run.font.size else None
                fc = run.font.color.rgb if run.font.color and run.font.color.type is not None else None
                runs.append({"text": run.text[:200], "size_pt": fs,
                             "bold": bool(run.font.bold), "color": f"#{fc}" if fc else None,
                             "font": run.font.name})
            texts.append({"runs": runs})
        item["texts"] = texts
    item["fill"] = color_of(shape.fill) if hasattr(shape, "fill") else None
    if shape.shape_type is not None and "PICTURE" in str(shape.shape_type):
        try:
            img = shape.image
            item["image"] = {"ext": img.ext, "size_px": [img.size[0], img.size[1]], "bytes": len(img.blob)}
        except Exception:
            item["image"] = None
    item["children"] = []
    if shape.shape_type is not None and "GROUP" in str(shape.shape_type):
        for sub in shape.shapes:
            shape_dump(sub, depth + 1, item["children"])
    out.append(item)


def main() -> int:
    ap = argparse.ArgumentParser()
    ap.add_argument("--file", required=True)
    ap.add_argument("--max-shapes", type=int, default=200)
    ap.add_argument("--slide", type=int)
    args = ap.parse_args()

    try:
        from pptx import Presentation

        prs = Presentation(args.file)
        slides_out = []
        slides = list(prs.slides)
        if args.slide is not None:
            slides = [slides[args.slide - 1]]
        total = 0
        truncated = False
        for idx, slide in enumerate(slides, start=1):
            shapes = []
            for shape in slide.shapes:
                if total >= args.max_shapes:
                    truncated = True
                    break
                shape_dump(shape, 0, shapes)
                total += 1
            slides_out.append({"slide": idx, "shapes": shapes, "shape_count": len(shapes)})
        print(S.dump_json({
            "schema": "vision-report/v2", "task": "pptx", "sensors": ["pptx"],
            "coordsys": "pt",
            "source": {"type": "pptx", "path": args.file,
                       "slide_size_pt": [(prs.slide_width or 0) / 12700.0, (prs.slide_height or 0) / 12700.0],
                       "total_slides": len(prs.slides)},
            "slides": slides_out, "truncated": truncated}))
        return 0
    except Exception as e:
        print(json.dumps({"error": "vs_pptx failed", "detail": str(e)[:500]}, ensure_ascii=False))
        return 1


if __name__ == "__main__":
    sys.exit(main())
