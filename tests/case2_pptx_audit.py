#!/usr/bin/env python3
"""Phase 2 用例2 — PPT 结构审计（DeepSeek 仅凭数字定位问题）。

生成含已知缺陷的夹具 → pptx_dump 无损导出 → 从数字计算重叠/对齐/对比度/出界。

DeepSeek 推理演示：审计规则全部基于 pptx_dump 输出的纯数字（pt 坐标、hex 填充、
字号），与模型所见数据完全一致 —— 见运行输出中的「审计表」与报告中的逐条定位。

运行: /home/Arch/conda-envs/pi-vision/bin/python tests/case2_pptx_audit.py
"""
import json
import subprocess
import sys
from pathlib import Path

ROOT = Path(__file__).resolve().parent.parent
PY = ROOT / "python"
OUT = ROOT / "tests" / "_out" / "case2"
PYBIN = "/home/Arch/conda-envs/pi-vision/bin/python"

PASS = 0
FAIL = 0


def check(name: str, cond: bool, detail: str = "") -> None:
    global PASS, FAIL
    if cond:
        PASS += 1
        print(f"  ✓ {name}")
    else:
        FAIL += 1
        print(f"  ✗ {name} {detail}")


def srgb_to_linear(c: float) -> float:
    c = c / 255.0
    return c / 12.92 if c <= 0.03928 else ((c + 0.055) / 1.055) ** 2.4


def contrast_ratio(fg_hex: str, bg_hex: str) -> float:
    def lum(h: str) -> float:
        rgb = tuple(int(h[i : i + 2], 16) for i in (1, 3, 5))
        r, g, b = (srgb_to_linear(c) for c in rgb)
        return 0.2126 * r + 0.7152 * g + 0.0722 * b

    l1, l2 = lum(fg_hex), lum(bg_hex)
    if l1 < l2:
        l1, l2 = l2, l1
    return (l1 + 0.05) / (l2 + 0.05)


def build_fixture() -> Path:
    """夹具：第1张含 4 类注入缺陷；第2张为整洁对照。4:3 = 720×540pt。"""
    from pptx import Presentation  # type: ignore[import-not-found]
    from pptx.dml.color import RGBColor  # type: ignore[import-not-found]
    from pptx.enum.shapes import MSO_SHAPE  # type: ignore[import-not-found]
    from pptx.util import Inches, Pt  # type: ignore[import-not-found]

    prs = Presentation()
    prs.slide_width = Inches(10)  # 720pt
    prs.slide_height = Inches(7.5)  # 540pt
    blank = prs.slide_layouts[6]

    def add_rect(slide, name, x, y, w, h, fill, text=None, font_pt: int = 14, font_color: str = "FFFFFF"):
        shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x / 72), Inches(y / 72),
                                     Inches(w / 72), Inches(h / 72))
        shp.name = name
        shp.fill.solid()
        shp.fill.fore_color.rgb = RGBColor.from_string(fill)
        shp.line.fill.background()
        if text is not None:
            tf = shp.text_frame
            tf.text = text
            tf.paragraphs[0].runs[0].font.size = Pt(font_pt)
            tf.paragraphs[0].runs[0].font.color.rgb = RGBColor.from_string(font_color)
        return shp

    s1 = prs.slides.add_slide(blank)
    t = s1.shapes.add_textbox(Inches(1), Inches(0.5), Inches(8), Inches(0.75))
    t.text_frame.text = "审计夹具：已知缺陷"
    t.text_frame.paragraphs[0].runs[0].font.size = Pt(24)
    t.text_frame.paragraphs[0].runs[0].font.color.rgb = RGBColor.from_string("1A1A1A")
    add_rect(s1, "cardA", 72, 144, 288, 180, "2278D2", "性能指标 98%", 14, "FFFFFF")
    add_rect(s1, "cardB", 288, 216, 288, 180, "FF6B6B")          # 与 A 重叠 72×108pt
    add_rect(s1, "cardC", 72, 432, 288, 90, "EAF2FB")            # 左缘与 A 对齐(72)
    add_rect(s1, "cardD", 120, 396, 216, 36, "F5F6F7")           # 左缘 120 ≠ 72（网格违规）
    add_rect(s1, "offSlide", 700, 480, 120, 48, "FFD166")        # 右缘 820 > 720（出界 100pt）

    s2 = prs.slides.add_slide(blank)
    t2 = s2.shapes.add_textbox(Inches(1), Inches(0.5), Inches(8), Inches(0.75))
    t2.text_frame.text = "整洁对照页"
    t2.text_frame.paragraphs[0].runs[0].font.size = Pt(24)
    t2.text_frame.paragraphs[0].runs[0].font.color.rgb = RGBColor.from_string("1A1A1A")
    add_rect(s2, "base1", 72, 144, 288, 180, "1F4E8C", "对齐基准", 14, "FFFFFF")
    add_rect(s2, "base2", 360, 144, 288, 180, "1F4E8C", "基准二", 14, "FFFFFF")

    f = OUT / "audit.pptx"
    prs.save(str(f))
    return f


def audit(d: dict) -> dict:
    """纯数字审计（与 DeepSeek 所见数据相同）。返回 findings 列表。"""
    W, H = d["source"]["slide_size_pt"]
    findings: list[dict] = []
    for slide in d["slides"]:
        shapes = slide["shapes"]
        idx = slide["slide"]
        for sh in shapes:
            x1, y1 = sh["pos_pt"]
            w, h = sh["size_pt"]
            x2, y2 = x1 + w, y1 + h
            sh["_x2"], sh["_y2"] = x2, y2
            if x2 > W + 0.01 or y2 > H + 0.01 or x1 < -0.01 or y1 < -0.01:
                overflow = max(x2 - W, y2 - H, -x1, -y1)
                findings.append({"slide": idx, "type": "off_slide", "shape": sh["name"],
                                 "overflow_pt": round(overflow, 1),
                                 "detail": f"右缘/下缘超出画布 {overflow:.1f}pt",
                                 "bbox": [x1, y1, x2, y2]})
        # 重叠：两两 bbox 交集
        for i in range(len(shapes)):
            a = shapes[i]
            for b in shapes[i + 1 :]:
                ix1 = max(a["pos_pt"][0], b["pos_pt"][0])
                iy1 = max(a["pos_pt"][1], b["pos_pt"][1])
                ix2 = min(a["_x2"], b["_x2"])
                iy2 = min(a["_y2"], b["_y2"])
                if ix2 > ix1 and iy2 > iy1:
                    area = (ix2 - ix1) * (iy2 - iy1)
                    findings.append({"slide": idx, "type": "overlap", "shapes": [a["name"], b["name"]],
                                     "intersect": [round(ix1, 1), round(iy1, 1), round(ix2, 1), round(iy2, 1)],
                                     "area_pt2": round(area, 1)})
        # 对齐：同排（y 区间重叠 ≥ 50% 较短者）且顶部相差 > 0.5pt
        for i in range(len(shapes)):
            a = shapes[i]
            for b in shapes[i + 1 :]:
                oy = min(a["_y2"], b["_y2"]) - max(a["pos_pt"][1], b["pos_pt"][1])
                min_h = min(a["size_pt"][1], b["size_pt"][1])
                if oy > 0.5 * min_h and abs(a["pos_pt"][1] - b["pos_pt"][1]) > 0.5:
                    findings.append({"slide": idx, "type": "top_misalign", "shapes": [a["name"], b["name"]],
                                     "top_gap_pt": round(abs(a["pos_pt"][1] - b["pos_pt"][1]), 1)})
        # 对比度：文本色 vs 形状填充
        for sh in shapes:
            fill = sh.get("fill")
            for para in sh.get("texts", []):
                for run in para.get("runs", []):
                    color = run.get("color")
                    if fill and fill.startswith("#") and color:
                        ratio = contrast_ratio(color, fill)
                        if ratio < 4.5:
                            findings.append({"slide": idx, "type": "contrast", "shape": sh["name"],
                                             "fg": color, "bg": fill, "ratio": round(ratio, 2),
                                             "size_pt": run.get("size_pt")})
    return {"findings": findings}


def main() -> int:
    OUT.mkdir(parents=True, exist_ok=True)
    fixture = build_fixture()

    d = run_pptx_dump(fixture)
    audit_result = audit(d)

    # 审计表（数字直出，供 DeepSeek 推理）
    print("[case2] pptx_dump 数字审计表")
    for slide in d["slides"]:
        print(f"  slide {slide['slide']}: 画布 {d['source']['slide_size_pt']}")
        for sh in slide["shapes"]:
            print(f"    {sh['name']:<10} pos={sh['pos_pt']} size={sh['size_pt']} "
                  f"fill={sh['fill']} texts={[r.get('color') for p in sh.get('texts', []) for r in p.get('runs', [])]}")
    print(f"  findings: {len(audit_result['findings'])}")
    for fnd in audit_result["findings"]:
        print(f"    {fnd}")

    print("[case2] 断言：注入缺陷全部被数字定位")
    s1_f = [f for f in audit_result["findings"] if f["slide"] == 1]
    s2_f = [f for f in audit_result["findings"] if f["slide"] == 2]

    ov = [f for f in s1_f if f["type"] == "overlap"]
    check("slide1 检出 cardA×cardB 重叠 72×108pt (7776pt²)",
          any(f["shapes"] == ["cardA", "cardB"] and abs(f["area_pt2"] - 7776) < 1 for f in ov),
          str(ov))
    check("slide1 无其他重叠", len(ov) == 1, str(ov))

    ct = [f for f in s1_f if f["type"] == "contrast"]
    check("slide1 检出 cardA 白字/蓝底对比度 ≈4.48 不达标",
          any(f["shape"] == "cardA" and 4.4 <= f["ratio"] <= 4.6 for f in ct), str(ct))

    al = [f for f in s1_f if f["type"] == "top_misalign"]
    check("slide1 检出 cardA×cardB 顶部错位 72pt",
          any(f["shapes"] == ["cardA", "cardB"] and f["top_gap_pt"] == 72.0 for f in al), str(al))

    off = [f for f in s1_f if f["type"] == "off_slide"]
    check("slide1 检出 offSlide 出界 100pt",
          any(f["shape"] == "offSlide" and abs(f["overflow_pt"] - 100) < 0.5 for f in off),
          str(off))

    check("slide2 整洁对照页零缺陷", len(s2_f) == 0, str(s2_f))

    # 对照：修复 cardA 文字颜色后对比度应达标（数字层面的修复验证）
    fixed_ratio = contrast_ratio("#FFFFFF", "#1F4E8C")
    check("修复参考：#FFFFFF on #1F4E8C = {:.1f}:1 达标 AA".format(fixed_ratio),
          fixed_ratio >= 4.5)

    print(f"\ncase2 结果: {PASS} 通过 / {FAIL} 失败")
    return 1 if FAIL else 0


def run_pptx_dump(fixture: Path) -> dict:
    r = subprocess.run([PYBIN, str(PY / "vs_pptx.py"), "--file", str(fixture)],
                       capture_output=True, text=True, timeout=120)
    if r.returncode != 0:
        raise AssertionError(f"pptx_dump exit {r.returncode}: {r.stderr[:300]}")
    try:
        return json.loads(r.stdout)
    except ValueError as e:
        raise AssertionError(f"bad JSON from pptx_dump: {r.stdout[:200]}") from e


if __name__ == "__main__":
    sys.exit(main())
