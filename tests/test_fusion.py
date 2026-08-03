#!/usr/bin/env python3
"""test_fusion.py — 融合层自测：crosscheck / audit / analyze 引擎。

运行（conda env）:
  /home/Arch/conda-envs/pi-vision/bin/python tests/test_fusion.py
"""
import json
import subprocess
import sys
from pathlib import Path

from PIL import Image, ImageDraw  # type: ignore[import-not-found]

ROOT = Path(__file__).resolve().parent.parent
PY = ROOT / "python"
OUT = ROOT / "tests" / "_out"
OUT.mkdir(parents=True, exist_ok=True)
PYBIN = "/home/Arch/conda-envs/pi-vision/bin/python"

PASS = 0
FAIL = 0


def run_tool(script: str, args: list[str]) -> dict:
    r = subprocess.run([PYBIN, str(PY / script), *args], capture_output=True, text=True, timeout=180)
    if r.returncode != 0:
        raise AssertionError(f"tool exit {r.returncode}: {r.stdout[:300]} {r.stderr[:200]}")
    try:
        return json.loads(r.stdout)
    except ValueError as e:
        raise AssertionError(f"bad JSON from {script}: {r.stdout[:200]}") from e


def check(name: str, cond: bool, detail: str = "") -> None:
    global PASS, FAIL
    if cond:
        PASS += 1
        print(f"  ✓ {name}")
    else:
        FAIL += 1
        print(f"  ✗ {name} {detail}")


def make_scene() -> tuple[Path, Path]:
    """合成场景：DOM 声明色 vs 实际渲染色不一致 + 一个 DOM 有文本但渲染为空白。"""
    from PIL import ImageFont  # type: ignore[import-not-found]

    try:
        f18 = ImageFont.truetype("/usr/share/fonts/noto-cjk/NotoSansCJK-Regular.ttc", 18)
    except OSError:
        f18 = ImageFont.load_default()
    W, H = 500, 300
    im = Image.new("RGB", (W, H), (255, 255, 255))
    d = ImageDraw.Draw(im)
    # 区域1：DOM 说 #111111（近黑），实际渲染成红色 → 颜色漂移
    d.rectangle([100, 50, 300, 80], fill=(255, 255, 255))
    d.text((100, 50), "标题测试", fill=(224, 0, 0), font=f18)
    # 区域2：DOM 有文本，渲染为空白 → 文本缺失（渲染失败）
    d.rectangle([100, 100, 300, 130], fill=(255, 255, 255))
    # 区域3：正常元素（DOM 与渲染一致），作对照
    d.rectangle([100, 160, 260, 190], fill=(255, 255, 255))
    d.text((100, 160), "正常元素", fill=(17, 17, 17), font=f18)
    img = OUT / "scene.png"
    im.save(img)

    dom = {
        "schema": "vision-report/v2", "task": "dom", "coordsys": "css_px",
        "source": {"dpr": 1.0},
        "elements": [
            {"id": 0, "type": "h1", "bbox": [100, 50, 300, 80], "text": "标题测试",
             "color": {"fill": "#FFFFFF", "text": "#111111"}, "source": ["dom"], "coordsys": "css_px"},
            {"id": 1, "type": "div", "bbox": [100, 100, 300, 130], "text": "canvas文本",
             "color": {"fill": "#FFFFFF", "text": "#222222"}, "source": ["dom"], "coordsys": "css_px"},
            {"id": 2, "type": "p", "bbox": [100, 160, 260, 190], "text": "正常元素",
             "color": {"fill": "#FFFFFF", "text": "#111111"}, "source": ["dom"], "coordsys": "css_px"},
        ],
    }
    dom_f = OUT / "scene_dom.json"
    dom_f.write_text(json.dumps(dom), encoding="utf-8")
    return img, dom_f


def test_crosscheck() -> None:
    print("[vs_crosscheck] 三方互验")
    img, dom_f = make_scene()
    ocr = run_tool("vs_ocr.py", ["--image", str(img), "--max-items", "20"])
    ocr_f = OUT / "scene_ocr.json"
    ocr_f.write_text(json.dumps(ocr), encoding="utf-8")

    d = run_tool("vs_crosscheck.py", ["--image", str(img), "--dom", str(dom_f),
                                      "--ocr", str(ocr_f), "--dpr", "1.0", "--color-threshold", "5.0"])
    types = [a["type"] for a in d["anomalies"]]
    drift = next((a for a in d["anomalies"] if a["type"] == "color_drift"), None)
    missing = next((a for a in d["anomalies"] if a["type"] == "text_missing_in_ocr"), None)
    check("检出颜色漂移（标题：声明 #111111 vs 实测红）",
          drift is not None and drift["evidence"]["dom_color"] == "#111111"
          and drift["evidence"]["delta_e76"] > 20, str(drift))
    check("检出文本缺失（canvas 区渲染为空白）",
          missing is not None and "canvas文本" in missing["evidence"]["dom_text"], str(missing))
    check("无多余噪音（anomaly ≤ 3）", len(d["anomalies"]) <= 3, str(types))


def test_audit() -> None:
    print("[vs_audit] 元素审计")
    from pptx import Presentation  # type: ignore[import-not-found]
    from pptx.dml.color import RGBColor
    from pptx.enum.shapes import MSO_SHAPE  # type: ignore[import-not-found]
    from pptx.util import Inches

    prs = Presentation()
    s = prs.slides.add_slide(prs.slide_layouts[6])
    a = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(1), Inches(1), Inches(3), Inches(2))
    a.fill.solid(); a.fill.fore_color.rgb = RGBColor(0x22, 0x78, 0xD2)
    b = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(3), Inches(2), Inches(2), Inches(2))
    b.fill.solid(); b.fill.fore_color.rgb = RGBColor(0xFF, 0x6B, 0x6B)
    off = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(-2), Inches(0.5), Inches(2), Inches(1))
    off.fill.solid(); off.fill.fore_color.rgb = RGBColor(0x44, 0x44, 0x44)
    tb = s.shapes.add_textbox(Inches(1), Inches(5), Inches(4), Inches(1))
    tb.text_frame.text = "白字浅底"
    tb.text_frame.paragraphs[0].runs[0].font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
    tb.fill.solid(); tb.fill.fore_color.rgb = RGBColor(0xFF, 0xFF, 0xEE)
    fx = OUT / "audit_fixture.pptx"
    prs.save(str(fx))

    pptx_report = run_tool("vs_pptx.py", ["--file", str(fx)])
    report_f = OUT / "audit_pptx.json"
    report_f.write_text(json.dumps(pptx_report), encoding="utf-8")
    d = run_tool("vs_audit.py", ["--report", str(report_f), "--canvas", "720x540"])
    types = {a["type"] for a in d["anomalies"]}
    check("检出元素重叠", "element_overlap" in types, str(types))
    check("检出出界元素（负坐标）", "off_canvas" in types, str(types))
    check("检出对比度不达标（白字浅底）", "contrast_fail" in types, str(types))


def test_analyze() -> None:
    print("[vs_analyze] 任务引擎（audit-pptx 配置）")
    d = run_tool("vs_analyze.py", ["--task", "audit-pptx",
                                   "--input", str(OUT / "audit_fixture.pptx")])
    check("报告为 schema v2 融合输出", d.get("schema") == "vision-report/v2" and "anomalies" in d, str(d)[:80])
    types = {a["type"] for a in d.get("anomalies", [])}
    check("引擎产出重叠/出界/对比度三类缺陷", {"element_overlap", "off_canvas", "contrast_fail"} <= types, str(types))


def main() -> int:
    print(f"fusion self-tests (python {sys.version.split()[0]})\n")
    test_crosscheck()
    test_audit()
    test_analyze()
    print(f"\n结果: {PASS} 通过 / {FAIL} 失败")
    return 1 if FAIL else 0


if __name__ == "__main__":
    sys.exit(main())
