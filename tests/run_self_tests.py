#!/usr/bin/env python3
"""pi-vision-struct 自测套件：真值断言，验证每个工具的确定性输出。

运行（conda env）:
  /home/Arch/conda-envs/pi-vision/bin/python tests/run_self_tests.py
"""
import json
import subprocess
import sys
from pathlib import Path

from PIL import Image, ImageDraw  # type: ignore[import-not-found]

ROOT = Path(__file__).resolve().parent.parent
PY = ROOT / "python"
OUT = ROOT / "tests" / "_out"
OUT.mkdir(parents=True, exist_ok=True)
PYBIN = "/home/Arch/conda-envs/pi-vision/bin/python"

PASS = 0
FAIL = 0


def run_tool(script: str, args: list[str]) -> dict:
    r = subprocess.run([PYBIN, str(PY / script), *args], capture_output=True, text=True, timeout=120)
    if r.returncode != 0:
        raise AssertionError(f"tool exit {r.returncode}: {r.stdout} {r.stderr[:300]}")
    try:
        return json.loads(r.stdout)
    except ValueError as e:
        raise AssertionError(f"bad JSON from {script}: {r.stdout[:200]}") from e


def check(name: str, cond: bool, detail: str = "") -> None:
    global PASS, FAIL
    if cond:
        PASS += 1
        print(f"  ✓ {name}")
    else:
        FAIL += 1
        print(f"  ✗ {name} {detail}")


def make_fixtures() -> Path:
    """Firefox 风格真值图（含 8px 微文字与已知色值）。"""
    W, H = 900, 600
    im = Image.new("RGB", (W, H), (245, 246, 247))
    d = ImageDraw.Draw(im)
    d.rectangle([0, 0, W, 38], fill=(50, 54, 60))
    d.rectangle([8, 156, 150, H - 8], fill=(34, 120, 210))  # 侧边栏 #2278D2
    d.rectangle([8, 104, W - 8, 150], fill=(253, 232, 230))
    d.text((20, 112), "页面加载失败  (错误代码: ERR_CONNECTION_REFUSED)", fill=(170, 30, 30))
    d.text((160, 396), "tiny:QWEN3VL8B-GT42-K7Z9", fill=(100, 100, 100))  # 8px 微文字
    p = OUT / "gt_firefox.png"
    im.save(p)
    return p


def test_pix(img: Path) -> None:
    print("[pix_analyze]")
    d = run_tool("vs_pix.py", ["--image", str(img), "--regions", "70,230,70,230", "14,20,14,20",
                               "--wcag", "2278D2,FFFFFF", "000000,FFFFFF", "--colors", "5"])
    colors = {c["hex"]: c["pct"] for c in d["metrics"]["dominant_colors"]}
    check("主背景 #F5F6F7 在直方图中", "#F5F6F7" in colors, str(colors))
    regs = {r["hex"] for r in d["regions"]}
    check("侧边栏区域色 == #2278D2", "#2278D2" in regs, str(regs))
    check("标题栏区域色 == #32363C", "#32363C" in regs, str(regs))
    w = d["wcag"][0]
    check("WCAG #2278D2 vs #FFF ratio≈4.48", 4.4 <= w["ratio"] <= 4.6, str(w))
    check("WCAG 判定与阈值一致 (4.48<4.5 不达标)", w["passes_aa"] == (w["ratio"] >= 4.5), str(w))
    w2 = d["wcag"][1]
    check("WCAG 黑 vs 白 = 21:1 达标 AA", abs(w2["ratio"] - 21.0) < 0.1 and w2["passes_aa"], str(w2))

    # diff 定位：注入异常后应定位到注入区域
    im = Image.open(img).convert("RGB")
    im2 = im.copy()
    ImageDraw.Draw(im2).rectangle([300, 110, 520, 145], fill=(120, 120, 120))
    bad = OUT / "gt_anomaly.png"
    im2.save(bad)
    d2 = run_tool("vs_pix.py", ["--image", str(img), "--compare", str(bad)])
    an = d2.get("anomalies", [])
    check("diff 定位到注入异常区", an and an[0]["bbox"][0] <= 302 <= an[0]["bbox"][2], str(an))


def test_ocr(img: Path) -> None:
    print("[ocr_boxes]")
    d = run_tool("vs_ocr.py", ["--image", str(img), "--max-items", "40"])
    texts = " ".join(it["text"] for it in d["elements"])
    check("识别微文字 QWEN3VL8B-GT42-K7Z9", "QWEN3VL8B-GT42-K7Z9" in texts.replace(" ", ""), texts[:80])
    check("识别错误码 ERR_CONNECTION_REFUSED", "ERR_CONNECTION_REFUSED" in texts.replace(" ", ""))
    check("坐标框为 4 点且中心在框内", all(
        it["bbox"][0] <= it["center"][0] <= it["bbox"][2] and it["bbox"][1] <= it["center"][1] <= it["bbox"][3]
        for it in d["elements"][:5]))


def test_pptx() -> None:
    print("[pptx_dump]")
    from pptx import Presentation  # type: ignore[import-not-found]
    from pptx.enum.shapes import MSO_SHAPE  # type: ignore[import-not-found]
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor

    prs = Presentation()
    s = prs.slides.add_slide(prs.slide_layouts[6])
    tb = s.shapes.add_textbox(Inches(1), Inches(0.5), Inches(8), Inches(1))
    tb.text_frame.text = "测试标题"
    tb.text_frame.paragraphs[0].runs[0].font.size = Pt(32)
    tb.text_frame.paragraphs[0].runs[0].font.color.rgb = RGBColor(0x22, 0x72, 0xD2)
    a = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(2), Inches(2), Inches(3), Inches(2))
    a.fill.solid(); a.fill.fore_color.rgb = RGBColor(0x22, 0x78, 0xD2)
    b = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(4), Inches(3), Inches(2), Inches(2))
    b.fill.solid(); b.fill.fore_color.rgb = RGBColor(0xFF, 0x6B, 0x6B)
    f = OUT / "fixture.pptx"
    prs.save(str(f))

    d = run_tool("vs_pptx.py", ["--file", str(f)])
    shapes = d["slides"][0]["shapes"]
    rects = [sh for sh in shapes if "AUTO_SHAPE" in str(sh["type"])]
    check("两个矩形都被导出", len(rects) == 2, str(len(rects)))
    a_rect = [sh for sh in rects if sh["fill"] == "#2278D2"][0]
    check("矩形A 位置 (144,144) 尺寸 216×144", a_rect["pos_pt"] == [144.0, 144.0] and a_rect["size_pt"] == [216.0, 144.0],
          str(a_rect["pos_pt"]) + str(a_rect["size_pt"]))
    title = [sh for sh in shapes if sh["type"].startswith("TEXT_BOX")][0]
    check("标题字体 32pt / 颜色 #2272D2",
          title["texts"][0]["runs"][0]["size_pt"] == 32.0 and title["texts"][0]["runs"][0]["color"] == "#2272D2",
          str(title["texts"]))


def test_dom() -> None:
    print("[dom_dump]")
    url = ('data:text/html,<h1 style="color:#2278D2;font-size:32px">修复测试</h1>'
           '<div style="position:absolute;left:100px;top:200px;width:300px;height:150px;'
           'background:#FF6B6B;z-index:5">占位</div>')
    d = run_tool("vs_dom.py", ["--url", url, "--max-elements", "10"])
    els = d["elements"]
    h1 = next((e for e in els if e["type"] == "text" and e["text"] == "修复测试"), None)
    div = next((e for e in els if e["type"] == "div"), None)
    check("h1 文本正确", h1 is not None and h1["text"] == "修复测试", str(h1))
    check("h1 颜色 == #2278D2", h1 is not None and (h1["color"] or {}).get("text") == "#2278D2", str(h1))
    check("div 绝对定位 z-index 5", div is not None and div["style"]["position"] == "absolute" and div["z"] == "5", str(div))
    check("div bbox = [100,200,400,350]", div is not None and div["bbox"] == [100, 200, 400, 350], str(div))


def main() -> int:
    print(f"pi-vision-struct self-tests (python {sys.version.split()[0]})\n")
    img = make_fixtures()
    test_pix(img)
    test_ocr(img)
    test_pptx()
    test_dom()
    print(f"\n结果: {PASS} 通过 / {FAIL} 失败")
    return 1 if FAIL else 0


if __name__ == "__main__":
    sys.exit(main())
